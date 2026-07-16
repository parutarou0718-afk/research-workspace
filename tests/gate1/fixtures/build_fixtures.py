from __future__ import annotations

import argparse
from datetime import datetime, timezone
import hashlib
from io import BytesIO
import json
from pathlib import Path
import tempfile
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from docx import Document
from docx.shared import Inches
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches as PptxInches
from pypdf import PdfWriter
from pypdf.generic import (
    ArrayObject,
    ByteStringObject,
    DecodedStreamObject,
    DictionaryObject,
    NameObject,
    NumberObject,
)


GENERATOR_VERSION = "gate1-docx-1"
FIXED_ZIP_TIME = (1980, 1, 1, 0, 0, 0)
FIXED_CORE_TIME = datetime(2000, 1, 1, tzinfo=timezone.utc)
ROOT = Path(__file__).resolve().parent
DOCX_ROOT = ROOT / "docx"
MANIFEST_PATH = ROOT / "manifest.json"

W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
REL = "http://schemas.openxmlformats.org/package/2006/relationships"
CT = "http://schemas.openxmlformats.org/package/2006/content-types"
WP = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
V = "urn:schemas-microsoft-com:vml"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
DCTERMS = "http://purl.org/dc/terms/"

PURPOSES = {
    "docx/body_order.docx": "OOXML body paragraph/table child order and heading path",
    "docx/table_escapes.docx": "versioned reversible table-cell escaping",
    "docx/image_alt.docx": "local image title and description alternative text",
    "docx/unsupported_constructs.docx": "stable warnings for excluded and unsupported OOXML",
    "pdf/normal_text.pdf": "two deterministic text pages",
    "pdf/empty_password.pdf": "empty-password encrypted text document",
    "pdf/password_required.pdf": "non-empty-password boundary fixture",
    "pdf/image_only.pdf": "valid raster-only page with no text-showing operator",
    "pdf/corrupt.pdf": "invalid non-PDF bytes",
    "pdf/truncated.pdf": "valid PDF with the final 32 bytes removed",
    "pptx/ordered_shapes.pptx": "slide and OOXML shape-tree order independent of coordinates",
    "pptx/nested_groups.pptx": "recursive shape paths and group depth safety",
    "pptx/table_alt_empty.pptx": "table escaping, explicit image alt, and an empty slide",
    "pptx/unsupported_constructs.pptx": "stable warnings for chart, SmartArt, media, embedded object, missing alt, and notes",
}

PDF_GENERATOR_VERSION = "gate1-pdf-1"
PPTX_GENERATOR_VERSION = "gate1-pptx-1"


def _new_document(title: str) -> Document:
    document = Document()
    properties = document.core_properties
    properties.title = title
    properties.author = "Research Workspace Fixture Builder"
    properties.subject = "Deterministic Gate 1 DOCX fixture"
    properties.keywords = "gate1,docx,fixture"
    properties.created = FIXED_CORE_TIME
    properties.modified = FIXED_CORE_TIME
    return document


def _zip_bytes(source: Path, transforms: dict[str, bytes] | None = None) -> dict[str, bytes]:
    with ZipFile(source, "r") as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    if transforms:
        members.update(transforms)
    return members


def _assert_no_external_relationships(members: dict[str, bytes]) -> None:
    for name, payload in members.items():
        if not name.endswith(".rels"):
            continue
        root = etree.fromstring(payload)
        for relationship in root.findall(f"{{{REL}}}Relationship"):
            if relationship.get("TargetMode") == "External":
                raise ValueError(f"external relationship forbidden in fixture: {name}")


def _write_canonical_package(destination: Path, members: dict[str, bytes]) -> None:
    settings_name = "word/settings.xml"
    if settings_name in members:
        settings = etree.fromstring(members[settings_name])
        zoom = settings.find(f"{{{W}}}zoom")
        if zoom is not None and zoom.get(f"{{{W}}}percent") is None:
            zoom.set(f"{{{W}}}percent", "100")
            members = {
                **members,
                settings_name: etree.tostring(
                    settings, xml_declaration=True, encoding="UTF-8", standalone=True
                ),
            }
    _assert_no_external_relationships(members)
    destination.parent.mkdir(parents=True, exist_ok=True)
    with ZipFile(destination, "w", compression=ZIP_DEFLATED, compresslevel=9) as archive:
        for name in sorted(members):
            info = ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = 0
            archive.writestr(info, members[name])


def _canonical_embedded_workbook(payload: bytes) -> bytes:
    with ZipFile(BytesIO(payload), "r") as archive:
        members = {name: archive.read(name) for name in archive.namelist()}
    core = etree.fromstring(members["docProps/core.xml"])
    for local_name in ("created", "modified"):
        element = core.find(f"{{{DCTERMS}}}{local_name}")
        assert element is not None
        element.text = "2000-01-01T00:00:00Z"
    members["docProps/core.xml"] = etree.tostring(
        core, xml_declaration=True, encoding="UTF-8", standalone=True
    )
    output = BytesIO()
    with ZipFile(output, "w", compression=ZIP_DEFLATED, compresslevel=9) as archive:
        for name in sorted(members):
            info = ZipInfo(name, FIXED_ZIP_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = 0
            archive.writestr(info, members[name])
    return output.getvalue()


def _save_canonical(document: Document, destination: Path, transform=None) -> None:
    with tempfile.TemporaryDirectory() as temporary_directory:
        raw = Path(temporary_directory) / "raw.docx"
        document.save(raw)
        members = _zip_bytes(raw)
    if transform is not None:
        members = transform(members)
    _write_canonical_package(destination, members)


def _body_order(destination: Path) -> None:
    document = _new_document("Body order fixture")
    document.add_heading("Methods", level=1)
    document.add_paragraph("P1")
    table = document.add_table(rows=2, cols=2)
    for cell, value in zip(
        (cell for row in table.rows for cell in row.cells),
        ("A", "B", "C", "D"),
        strict=True,
    ):
        cell.text = value
    document.add_paragraph("P2")
    _save_canonical(document, destination)


def _table_escapes(destination: Path) -> None:
    document = _new_document("Table escaping fixture")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "back\\slash"
    table.cell(0, 1).text = "tab\tvalue"
    table.cell(1, 0).text = "line\nfeed"
    table.cell(1, 1).text = "CARRIAGE_MARKER"

    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        xml = members["word/document.xml"]
        xml = xml.replace(b"CARRIAGE_MARKER", b"carriage&#13;return")
        return {**members, "word/document.xml": xml}

    _save_canonical(document, destination, transform)


def _image_alt(destination: Path) -> None:
    document = _new_document("Image alt fixture")
    with tempfile.TemporaryDirectory() as temporary_directory:
        image_path = Path(temporary_directory) / "fixture.png"
        Image.new("RGB", (10, 10), (32, 96, 160)).save(
            image_path, format="PNG", optimize=False, compress_level=9
        )
        document.add_picture(str(image_path), width=Inches(0.25), height=Inches(0.25))

    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        root = etree.fromstring(members["word/document.xml"])
        doc_properties = root.xpath("//wp:docPr", namespaces={"wp": WP})
        if len(doc_properties) != 1:
            raise AssertionError("image fixture must contain one wp:docPr")
        doc_properties[0].set("title", "Figure title")
        doc_properties[0].set("descr", "Figure description")
        return {
            **members,
            "word/document.xml": etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
        }

    _save_canonical(document, destination, transform)


def _unsupported_constructs(destination: Path) -> None:
    document = _new_document("Unsupported constructs fixture")
    document.add_paragraph("Supported body")
    document.sections[0].header.paragraphs[0].text = "Excluded header"
    document.sections[0].footer.paragraphs[0].text = "Excluded footer"

    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        root = etree.fromstring(members["word/document.xml"])
        body = root.find(f"{{{W}}}body")
        assert body is not None
        section_properties = body.find(f"{{{W}}}sectPr")

        def insert(element: etree._Element) -> None:
            index = len(body) if section_properties is None else body.index(section_properties)
            body.insert(index, element)

        textbox_paragraph = etree.Element(f"{{{W}}}p")
        textbox_run = etree.SubElement(textbox_paragraph, f"{{{W}}}r")
        picture = etree.SubElement(textbox_run, f"{{{W}}}pict")
        shape = etree.SubElement(picture, f"{{{V}}}shape")
        shape.set("id", "Gate1TextBox")
        shape.set("style", "width:100pt;height:20pt")
        vml_textbox = etree.SubElement(shape, f"{{{V}}}textbox")
        textbox_content = etree.SubElement(vml_textbox, f"{{{W}}}txbxContent")
        nested_paragraph = etree.SubElement(textbox_content, f"{{{W}}}p")
        nested_run = etree.SubElement(nested_paragraph, f"{{{W}}}r")
        etree.SubElement(nested_run, f"{{{W}}}t").text = "Unsupported text box"
        insert(textbox_paragraph)

        field_paragraph = etree.Element(f"{{{W}}}p")
        field = etree.SubElement(field_paragraph, f"{{{W}}}fldSimple")
        field.set(f"{{{W}}}instr", " DATE ")
        field_run = etree.SubElement(field, f"{{{W}}}r")
        etree.SubElement(field_run, f"{{{W}}}t").text = "Field result"
        insert(field_paragraph)

        equation_paragraph = etree.Element(f"{{{W}}}p")
        equation = etree.SubElement(equation_paragraph, f"{{{M}}}oMath")
        equation_run = etree.SubElement(equation, f"{{{M}}}r")
        etree.SubElement(equation_run, f"{{{M}}}t").text = "x+y"
        insert(equation_paragraph)

        tracked_paragraph = etree.Element(f"{{{W}}}p")
        insertion = etree.SubElement(tracked_paragraph, f"{{{W}}}ins")
        insertion.set(f"{{{W}}}id", "1")
        insertion.set(f"{{{W}}}author", "Research Workspace Fixture Builder")
        insertion.set(f"{{{W}}}date", "2000-01-01T00:00:00Z")
        tracked_run = etree.SubElement(insertion, f"{{{W}}}r")
        etree.SubElement(tracked_run, f"{{{W}}}t").text = "Tracked insertion"
        insert(tracked_paragraph)

        object_paragraph = etree.Element(f"{{{W}}}p")
        object_run = etree.SubElement(object_paragraph, f"{{{W}}}r")
        etree.SubElement(object_run, f"{{{W}}}object")
        insert(object_paragraph)

        document_xml = etree.tostring(
            root, xml_declaration=True, encoding="UTF-8", standalone=True
        )

        relationships = etree.fromstring(members["word/_rels/document.xml.rels"])
        existing_ids = {
            int(value[3:])
            for value in (item.get("Id", "") for item in relationships)
            if value.startswith("rId") and value[3:].isdigit()
        }
        next_id = max(existing_ids, default=0) + 1

        def relationship(rel_type: str, target: str) -> None:
            nonlocal next_id
            item = etree.SubElement(relationships, f"{{{REL}}}Relationship")
            item.set("Id", f"rId{next_id}")
            item.set("Type", rel_type)
            item.set("Target", target)
            next_id += 1

        relationship(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/footnotes",
            "footnotes.xml",
        )
        relationship(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/endnotes",
            "endnotes.xml",
        )
        relationship(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments",
            "comments.xml",
        )
        relationship(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData",
            "diagrams/data1.xml",
        )
        relationship(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
            "embeddings/oleObject1.bin",
        )

        content_types = etree.fromstring(members["[Content_Types].xml"])

        def override(part_name: str, content_type: str) -> None:
            item = etree.SubElement(content_types, f"{{{CT}}}Override")
            item.set("PartName", part_name)
            item.set("ContentType", content_type)

        override(
            "/word/footnotes.xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.footnotes+xml",
        )
        override(
            "/word/endnotes.xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.endnotes+xml",
        )
        override(
            "/word/comments.xml",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml",
        )
        override(
            "/word/diagrams/data1.xml",
            "application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml",
        )
        override(
            "/word/embeddings/oleObject1.bin",
            "application/vnd.openxmlformats-officedocument.oleObject",
        )

        return {
            **members,
            "word/document.xml": document_xml,
            "word/_rels/document.xml.rels": etree.tostring(
                relationships, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
            "[Content_Types].xml": etree.tostring(
                content_types, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
            "word/footnotes.xml": (
                f'<w:footnotes xmlns:w="{W}"><w:footnote w:id="1"><w:p><w:r>'
                '<w:t>Footnote</w:t></w:r></w:p></w:footnote></w:footnotes>'
            ).encode(),
            "word/endnotes.xml": (
                f'<w:endnotes xmlns:w="{W}"><w:endnote w:id="1"><w:p><w:r>'
                '<w:t>Endnote</w:t></w:r></w:p></w:endnote></w:endnotes>'
            ).encode(),
            "word/comments.xml": (
                f'<w:comments xmlns:w="{W}"><w:comment w:id="0" '
                'w:author="Research Workspace Fixture Builder" w:date="2000-01-01T00:00:00Z"><w:p><w:r>'
                '<w:t>Comment</w:t></w:r></w:p></w:comment></w:comments>'
            ).encode(),
            "word/diagrams/data1.xml": b"<dgm:dataModel xmlns:dgm=\"http://schemas.openxmlformats.org/drawingml/2006/diagram\"/>",
            "word/embeddings/oleObject1.bin": b"deterministic-ole-fixture",
        }

    _save_canonical(document, destination, transform)


def _pdf_writer(title: str) -> PdfWriter:
    writer = PdfWriter()
    binary_marker = DecodedStreamObject()
    binary_marker.set_data(b"\x00gate1-binary-pdf-marker")
    writer._add_object(binary_marker)
    writer.add_metadata(
        {
            "/Title": title,
            "/Author": "Research Workspace Fixture Builder",
            "/Subject": "Deterministic Gate 1 PDF fixture",
            "/Keywords": "gate1,pdf,fixture",
            "/CreationDate": "D:20000101000000Z",
            "/ModDate": "D:20000101000000Z",
        }
    )
    identifier = ByteStringObject(b"gate1-pdf-id-0001")
    writer._ID = ArrayObject((identifier, identifier))
    return writer


def _add_text_page(writer: PdfWriter, text: str) -> None:
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_reference = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): font_reference}
            )
        }
    )
    content = DecodedStreamObject()
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    content.set_data(f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii"))
    page[NameObject("/Contents")] = writer._add_object(content)


def _add_image_only_page(writer: PdfWriter) -> None:
    page = writer.add_blank_page(width=612, height=792)
    image = DecodedStreamObject()
    image.set_data(bytes((32, 96, 160)) * 100)
    image.update(
        {
            NameObject("/Type"): NameObject("/XObject"),
            NameObject("/Subtype"): NameObject("/Image"),
            NameObject("/Width"): NumberObject(10),
            NameObject("/Height"): NumberObject(10),
            NameObject("/ColorSpace"): NameObject("/DeviceRGB"),
            NameObject("/BitsPerComponent"): NumberObject(8),
        }
    )
    image_reference = writer._add_object(image)
    page[NameObject("/Resources")] = DictionaryObject(
        {
            NameObject("/XObject"): DictionaryObject(
                {NameObject("/Im0"): image_reference}
            )
        }
    )
    content = DecodedStreamObject()
    content.set_data(b"q 10 0 0 10 72 700 cm /Im0 Do Q")
    page[NameObject("/Contents")] = writer._add_object(content)


def _pdf_bytes(
    title: str,
    texts: tuple[str, ...] = (),
    *,
    image_only: bool = False,
    encryption_password: str | None = None,
) -> bytes:
    writer = _pdf_writer(title)
    for text in texts:
        _add_text_page(writer, text)
    if image_only:
        _add_image_only_page(writer)
    if encryption_password is not None:
        writer.encrypt(encryption_password, algorithm="RC4-128")
    output = BytesIO()
    writer.write(output)
    payload = output.getvalue()
    if encryption_password is not None:
        payload += b"\n%\x00gate1-encrypted-binary-marker\n"
    return payload


def _build_pdf_fixtures(output_root: Path) -> None:
    pdf_root = output_root / "pdf"
    pdf_root.mkdir(parents=True, exist_ok=True)
    normal = _pdf_bytes("Normal text fixture", ("Page one text", "Page two text"))
    (pdf_root / "normal_text.pdf").write_bytes(normal)
    (pdf_root / "empty_password.pdf").write_bytes(
        _pdf_bytes(
            "Empty password fixture",
            ("Page one text", "Page two text"),
            encryption_password="",
        )
    )
    (pdf_root / "password_required.pdf").write_bytes(
        _pdf_bytes(
            "Password required fixture",
            ("Protected page text",),
            encryption_password="gate1-fixture-password",
        )
    )
    (pdf_root / "image_only.pdf").write_bytes(
        _pdf_bytes("Image only fixture", image_only=True)
    )
    (pdf_root / "corrupt.pdf").write_bytes(
        b"not-a-pdf\x00deterministic-corrupt-fixture\n"
    )
    (pdf_root / "truncated.pdf").write_bytes(normal[:-32])


def _new_presentation(title: str) -> Presentation:
    presentation = Presentation()
    properties = presentation.core_properties
    properties.title = title
    properties.author = "Research Workspace Fixture Builder"
    properties.subject = "Deterministic Gate 1 PPTX fixture"
    properties.keywords = "gate1,pptx,fixture"
    properties.created = FIXED_CORE_TIME
    properties.modified = FIXED_CORE_TIME
    return presentation


def _save_canonical_presentation(
    presentation: Presentation,
    destination: Path,
    transform=None,
) -> None:
    with tempfile.TemporaryDirectory() as temporary_directory:
        raw = Path(temporary_directory) / "raw.pptx"
        presentation.save(raw)
        members = _zip_bytes(raw)
    if transform is not None:
        members = transform(members)
    _write_canonical_package(destination, members)


def _ordered_shapes_pptx(destination: Path) -> None:
    presentation = _new_presentation("Ordered shapes fixture")
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    first = slide.shapes.add_textbox(PptxInches(6), PptxInches(4), PptxInches(2), PptxInches(1))
    first.text = "Slide 1 XML first"
    second = slide.shapes.add_textbox(PptxInches(1), PptxInches(1), PptxInches(2), PptxInches(1))
    second.text = "Slide 1 XML second"
    slide = presentation.slides.add_slide(blank)
    first = slide.shapes.add_textbox(PptxInches(5), PptxInches(3), PptxInches(2), PptxInches(1))
    first.text = "Slide 2 XML first"
    second = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.5), PptxInches(2), PptxInches(1))
    second.text = "Slide 2 XML second"
    _save_canonical_presentation(presentation, destination)


def _nested_groups_pptx(destination: Path) -> None:
    presentation = _new_presentation("Nested groups fixture")
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    outer = slide.shapes.add_group_shape()
    outer_leaf = outer.shapes.add_textbox(0, 0, PptxInches(2), PptxInches(0.5))
    outer_leaf.text = "Outer leaf"
    inner = outer.shapes.add_group_shape()
    nested_leaf = inner.shapes.add_textbox(0, 0, PptxInches(2), PptxInches(0.5))
    nested_leaf.text = "Nested leaf"

    deep_slide = presentation.slides.add_slide(blank)
    group = deep_slide.shapes.add_group_shape()
    for _ in range(32):
        group = group.shapes.add_group_shape()
    too_deep = group.shapes.add_textbox(0, 0, PptxInches(2), PptxInches(0.5))
    too_deep.text = "Too deep leaf"

    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        root = etree.fromstring(members["ppt/slides/slide1.xml"])
        shape_tree = root.find(f".//{{{P}}}spTree")
        assert shape_tree is not None
        outer_group = shape_tree.findall(f"{{{P}}}grpSp")[0]
        outer_group.find(f"{{{P}}}nvGrpSpPr/{{{P}}}cNvPr").set("id", "2")
        inner_group = outer_group.findall(f"{{{P}}}grpSp")[0]
        inner_group.find(f"{{{P}}}nvGrpSpPr/{{{P}}}cNvPr").set("id", "5")
        inner_shape = inner_group.findall(f"{{{P}}}sp")[0]
        inner_shape.find(f"{{{P}}}nvSpPr/{{{P}}}cNvPr").set("id", "9")
        return {
            **members,
            "ppt/slides/slide1.xml": etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
        }

    _save_canonical_presentation(presentation, destination, transform)


def _table_alt_empty_pptx(destination: Path) -> None:
    presentation = _new_presentation("Table alt empty fixture")
    blank = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank)
    table = slide.shapes.add_table(2, 2, PptxInches(0.5), PptxInches(0.5), PptxInches(5), PptxInches(2)).table
    table.cell(0, 0).text = "back\\slash"
    table.cell(0, 1).text = "TAB_MARKER"
    table.cell(1, 0).text = "line\nfeed"
    table.cell(1, 1).text = "CR_MARKER"
    with tempfile.TemporaryDirectory() as temporary_directory:
        image_path = Path(temporary_directory) / "fixture.png"
        Image.new("RGB", (10, 10), (32, 96, 160)).save(
            image_path, format="PNG", optimize=False, compress_level=9
        )
        slide.shapes.add_picture(
            str(image_path), PptxInches(6), PptxInches(1), PptxInches(0.5), PptxInches(0.5)
        )
    presentation.slides.add_slide(blank)

    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        root = etree.fromstring(members["ppt/slides/slide1.xml"])
        text_nodes = root.findall(f".//{{{A}}}t")
        for node in text_nodes:
            if node.text == "TAB_MARKER":
                node.text = "tab\tvalue"
            elif node.text == "CR_MARKER":
                node.text = "carriage\rreturn"
        pictures = root.findall(f".//{{{P}}}pic")
        assert len(pictures) == 1
        properties = pictures[0].find(f"{{{P}}}nvPicPr/{{{P}}}cNvPr")
        assert properties is not None
        properties.set("title", "Figure title")
        properties.set("descr", "Figure description")
        return {
            **members,
            "ppt/slides/slide1.xml": etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
        }

    _save_canonical_presentation(presentation, destination, transform)


def _unsupported_constructs_pptx(destination: Path) -> None:
    presentation = _new_presentation("Unsupported constructs fixture")
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ("One", "Two")
    chart_data.add_series("Fixture", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PptxInches(0.5),
        PptxInches(0.5),
        PptxInches(3),
        PptxInches(2),
        chart_data,
    )
    with tempfile.TemporaryDirectory() as temporary_directory:
        image_path = Path(temporary_directory) / "missing-alt.png"
        Image.new("RGB", (10, 10), (160, 96, 32)).save(
            image_path, format="PNG", optimize=False, compress_level=9
        )
        slide.shapes.add_picture(
            str(image_path), PptxInches(4), PptxInches(0.5), PptxInches(0.5), PptxInches(0.5)
        )
    slide.notes_slide.notes_text_frame.text = "Speaker notes"

    def transform(members: dict[str, bytes]) -> dict[str, bytes]:
        workbook_name = "ppt/embeddings/Microsoft_Excel_Sheet1.xlsx"
        members = {
            **members,
            workbook_name: _canonical_embedded_workbook(members[workbook_name]),
        }
        root = etree.fromstring(members["ppt/slides/slide1.xml"])
        shape_tree = root.find(f".//{{{P}}}spTree")
        assert shape_tree is not None
        picture_properties = shape_tree.find(f".//{{{P}}}pic/{{{P}}}nvPicPr/{{{P}}}cNvPr")
        assert picture_properties is not None
        picture_properties.attrib.pop("title", None)
        picture_properties.attrib.pop("descr", None)
        next_id = max(
            int(item.get("id", "0"))
            for item in shape_tree.findall(f".//{{{P}}}cNvPr")
        ) + 1
        def graphic_frame(name: str, uri: str, marker: str) -> None:
            nonlocal next_id
            frame = etree.fromstring(
                (
                    f'<p:graphicFrame xmlns:p="{P}" xmlns:a="{A}">'
                    f'<p:nvGraphicFramePr><p:cNvPr id="{next_id}" name="{name}"/>'
                    '<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>'
                    '<p:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></p:xfrm>'
                    f'<a:graphic><a:graphicData uri="{uri}"><a:{marker}/>'
                    '</a:graphicData></a:graphic></p:graphicFrame>'
                ).encode()
            )
            next_id += 1
            shape_tree.append(frame)

        graphic_frame(
            "SmartArt fixture",
            "http://schemas.openxmlformats.org/drawingml/2006/diagram",
            "ext",
        )
        graphic_frame(
            "Media fixture",
            "http://schemas.microsoft.com/office/powerpoint/2010/main/media",
            "ext",
        )
        graphic_frame(
            "Embedded object fixture",
            "http://schemas.openxmlformats.org/presentationml/2006/ole",
            "ext",
        )
        return {
            **members,
            "ppt/slides/slide1.xml": etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            ),
        }

    _save_canonical_presentation(presentation, destination, transform)


def _build_pptx_fixtures(output_root: Path) -> None:
    pptx_root = output_root / "pptx"
    builders = {
        "ordered_shapes.pptx": _ordered_shapes_pptx,
        "nested_groups.pptx": _nested_groups_pptx,
        "table_alt_empty.pptx": _table_alt_empty_pptx,
        "unsupported_constructs.pptx": _unsupported_constructs_pptx,
    }
    for name, builder in builders.items():
        builder(pptx_root / name)


def _build(output_root: Path) -> bytes:
    docx_root = output_root / "docx"
    builders = {
        "body_order.docx": _body_order,
        "table_escapes.docx": _table_escapes,
        "image_alt.docx": _image_alt,
        "unsupported_constructs.docx": _unsupported_constructs,
    }
    for name, builder in builders.items():
        builder(docx_root / name)
    _build_pdf_fixtures(output_root)
    _build_pptx_fixtures(output_root)

    entries = []
    for relative_path, purpose in PURPOSES.items():
        path = output_root / relative_path
        payload = path.read_bytes()
        entries.append(
            {
                "relative_path": relative_path,
                "sha256": hashlib.sha256(payload).hexdigest(),
                "size_bytes": len(payload),
                "purpose": purpose,
                "generator_version": (
                    PDF_GENERATOR_VERSION
                    if relative_path.startswith("pdf/")
                    else (
                        PPTX_GENERATOR_VERSION
                        if relative_path.startswith("pptx/")
                        else GENERATOR_VERSION
                    )
                ),
            }
        )
    manifest = {"schema_version": "1.0", "fixtures": entries}
    return (json.dumps(manifest, ensure_ascii=False, indent=2) + "\n").encode("utf-8")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--check", action="store_true")
    arguments = parser.parse_args()
    if arguments.check:
        with tempfile.TemporaryDirectory() as temporary_directory:
            generated_root = Path(temporary_directory)
            generated_manifest = _build(generated_root)
            if not MANIFEST_PATH.exists() or MANIFEST_PATH.read_bytes() != generated_manifest:
                raise SystemExit("fixture manifest is not reproducible")
            for relative_path in PURPOSES:
                if (ROOT / relative_path).read_bytes() != (generated_root / relative_path).read_bytes():
                    raise SystemExit(f"fixture is not reproducible: {relative_path}")
        return 0

    manifest = _build(ROOT)
    MANIFEST_PATH.write_bytes(manifest)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
