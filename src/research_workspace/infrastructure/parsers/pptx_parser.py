"""Deterministic python-pptx adapter over one declared immutable snapshot."""

from __future__ import annotations

from collections.abc import Mapping, Sequence
from importlib.metadata import version
from zipfile import BadZipFile

from lxml.etree import XMLSyntaxError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError

from research_workspace.application.dto.parsing_dto import ParseRequest, ParseResult
from research_workspace.domain.parsing import (
    PARSER_WARNING_CODES,
    ParseContractError,
    build_parse_artifact_identity,
    canonicalize_warnings,
    expand_parser_config,
    make_block_id,
    normalize_quote,
    validate_parsed_document_v2,
)
from research_workspace.infrastructure.parsers.table_text import (
    TABLE_TEXT_VERSION,
    escape_table_tsv,
)


_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_MAX_SHAPE_PATH = 32
_PARAGRAPH_KINDS = frozenset({"paragraph", "footnote", "image_alt"})


def _native_locator(
    slide: int,
    shape_path: Sequence[Mapping[str, int]],
    *,
    text_frame_paragraph_index: int | None = None,
    table_index: int | None = None,
    notes: bool = False,
    alt_source: str | None = None,
) -> dict[str, object]:
    return {
        "type": "pptx",
        "slide": slide,
        "shape_path": [dict(item) for item in shape_path],
        "text_frame_paragraph_index": text_frame_paragraph_index,
        "table_index": table_index,
        "row_index": None,
        "column_index": None,
        "notes": notes,
        "alt_source": alt_source,
    }


def _warning(code: str, native_locator: Mapping[str, object] | None = None) -> dict[str, object]:
    if code not in PARSER_WARNING_CODES:  # pragma: no cover - adapter registry invariant
        raise ValueError("warning is not registered")
    return {
        "code": code,
        "block_index": None,
        "native_locator": dict(native_locator) if native_locator is not None else None,
    }


class _DocumentBuilder:
    def __init__(self, artifact_id: object) -> None:
        self.artifact_id = artifact_id
        self.blocks: list[dict[str, object]] = []
        self.paragraph_index = 0

    def add(
        self,
        kind: str,
        text: str,
        native_locator: Mapping[str, object],
        metadata: Mapping[str, object],
        *,
        paragraph_like: bool,
    ) -> None:
        normalized = normalize_quote(text)
        if not normalized:
            return
        block_index = len(self.blocks)
        paragraph_index = self.paragraph_index if paragraph_like else None
        locator: dict[str, object] = {
            "page": None,
            "slide": native_locator["slide"],
            "block_index": block_index,
            "paragraph_index": paragraph_index,
            "paragraph_id": None,
            "heading_path": [],
            "char_start": 0,
            "char_end": len(normalized),
            "source_offset_start": None,
            "source_offset_end": None,
            "bbox": None,
            "native_locator": dict(native_locator),
        }
        block_id = make_block_id(self.artifact_id, block_index, kind, locator, normalized)
        if kind in _PARAGRAPH_KINDS:
            locator["paragraph_id"] = block_id
        self.blocks.append(
            {
                "block_id": block_id,
                "block_index": block_index,
                "kind": kind,
                "text": normalized,
                "locator": locator,
                "metadata": dict(metadata),
            }
        )
        if paragraph_like:
            self.paragraph_index += 1


class PptxParser:
    """Parse PPTX slide trees without executing or resolving linked content."""

    parser_id = "python-pptx"
    parser_version = version("python-pptx")
    supported_mime_types = frozenset({_PPTX_MIME})

    def parse(self, request: ParseRequest) -> ParseResult:
        try:
            config = expand_parser_config(request.parser_config)
        except ParseContractError:
            return ParseResult(None, (), "UNSUPPORTED_CONFIGURATION")
        if request.mime_type not in self.supported_mime_types or config["ocr_enabled"] is not False:
            return ParseResult(None, (), "UNSUPPORTED_CONFIGURATION")

        try:
            presentation = Presentation(request.snapshot_path)
        except (BadZipFile, PackageNotFoundError, FileNotFoundError):
            return ParseResult(None, (), "PPTX_INVALID_PACKAGE")
        except (KeyError, ValueError, XMLSyntaxError):
            return ParseResult(None, (), "PPTX_CORRUPT")
        except (PermissionError, OSError):
            return ParseResult(None, (), "PPTX_INVALID_PACKAGE")

        identity = build_parse_artifact_identity(
            request.snapshot_id,
            self.parser_id,
            self.parser_version,
            config,
            "2.0",
        )
        builder = _DocumentBuilder(request.parse_artifact_id)
        warnings: list[dict[str, object]] = [_warning("PPTX_READING_ORDER_HEURISTIC")]
        for slide_number, slide in enumerate(presentation.slides, start=1):
            table_index = 0
            for tree_index, shape in enumerate(slide.shapes):
                path = ({"shape_id": shape.shape_id, "tree_index": tree_index},)
                table_index = self._visit_shape(
                    shape,
                    slide_number,
                    path,
                    table_index,
                    builder,
                    warnings,
                    bool(config["preserve_image_alt"]),
                )
            self._visit_notes(
                slide,
                slide_number,
                builder,
                warnings,
                bool(config["preserve_footnotes"]),
            )

        canonical_warnings = canonicalize_warnings(warnings)
        keywords: list[str] = []
        for raw_keyword in (presentation.core_properties.keywords or "").split(","):
            keyword = raw_keyword.strip()
            if keyword and keyword not in keywords:
                keywords.append(keyword)
        parsed_document = {
            "schema_version": "2.0",
            "parse_artifact_id": str(request.parse_artifact_id),
            "source": {
                "source_snapshot_id": str(request.snapshot_id),
                "sha256": request.snapshot_sha256,
                "mime_type": request.mime_type,
                "size_bytes": request.snapshot_path.stat().st_size,
                "storage_relative_path": (
                    f"sources/sha256/{request.snapshot_sha256[:2]}/"
                    f"{request.snapshot_sha256}/content"
                ),
            },
            "parser": {
                "parser_id": self.parser_id,
                "parser_version": self.parser_version,
                "config_fingerprint": identity.config_fingerprint,
                "contract_version": "2.0",
            },
            "title": presentation.core_properties.title or None,
            "metadata": {
                "language": config["language"],
                "page_count": None,
                "slide_count": len(presentation.slides),
                "author": presentation.core_properties.author or None,
                "subject": presentation.core_properties.subject or None,
                "keywords": keywords,
            },
            "blocks": builder.blocks,
            "warnings": list(canonical_warnings),
        }
        try:
            validate_parsed_document_v2(parsed_document)
        except ParseContractError:
            return ParseResult(
                None,
                tuple(dict.fromkeys(item["code"] for item in canonical_warnings)),
                "PARSED_DOCUMENT_CONTRACT_INVALID",
            )
        return ParseResult(
            parsed_document,
            tuple(dict.fromkeys(item["code"] for item in canonical_warnings)),
            None,
        )

    def _visit_shape(
        self,
        shape,
        slide_number: int,
        shape_path: tuple[dict[str, int], ...],
        table_index: int,
        builder: _DocumentBuilder,
        warnings: list[dict[str, object]],
        preserve_image_alt: bool,
    ) -> int:
        native = _native_locator(slide_number, shape_path)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if len(shape_path) == _MAX_SHAPE_PATH and len(shape.shapes):
                warnings.append(_warning("PPTX_GROUP_DEPTH_EXCEEDED", native))
                return table_index
            for tree_index, child in enumerate(shape.shapes):
                child_path = shape_path + (
                    {"shape_id": child.shape_id, "tree_index": tree_index},
                )
                table_index = self._visit_shape(
                    child,
                    slide_number,
                    child_path,
                    table_index,
                    builder,
                    warnings,
                    preserve_image_alt,
                )
            return table_index

        if getattr(shape, "has_table", False):
            rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
            builder.add(
                "table",
                escape_table_tsv(rows),
                _native_locator(slide_number, shape_path, table_index=table_index),
                {
                    "table_text_version": TABLE_TEXT_VERSION,
                    "row_count": len(rows),
                    "column_count": max((len(row) for row in rows), default=0),
                },
                paragraph_like=False,
            )
            return table_index + 1

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            self._visit_picture(
                shape,
                slide_number,
                shape_path,
                builder,
                warnings,
                preserve_image_alt,
            )
            return table_index

        unsupported_code = self._unsupported_code(shape)
        if unsupported_code is not None:
            warnings.append(_warning(unsupported_code, native))
            return table_index

        if getattr(shape, "has_text_frame", False):
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                builder.add(
                    "paragraph",
                    paragraph.text,
                    _native_locator(
                        slide_number,
                        shape_path,
                        text_frame_paragraph_index=paragraph_index,
                    ),
                    {},
                    paragraph_like=True,
                )
        return table_index

    @staticmethod
    def _unsupported_code(shape) -> str | None:
        if getattr(shape, "has_chart", False):
            return "PPTX_CHART_TEXT_UNSUPPORTED"
        if shape.shape_type in {
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
        }:
            return "PPTX_EMBEDDED_OBJECT_UNSUPPORTED"
        if shape.shape_type in {MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.WEB_VIDEO}:
            return "PPTX_MEDIA_UNSUPPORTED"
        try:
            uri = shape._element.graphic.graphicData.uri
        except AttributeError:
            return None
        if uri.endswith("/diagram"):
            return "PPTX_SMARTART_UNSUPPORTED"
        if "media" in uri:
            return "PPTX_MEDIA_UNSUPPORTED"
        if uri.endswith("/ole"):
            return "PPTX_EMBEDDED_OBJECT_UNSUPPORTED"
        return None

    @staticmethod
    def _visit_picture(
        shape,
        slide_number: int,
        shape_path: tuple[dict[str, int], ...],
        builder: _DocumentBuilder,
        warnings: list[dict[str, object]],
        preserve_image_alt: bool,
    ) -> None:
        properties = shape._element.nvPicPr.cNvPr
        values = (
            ("title", properties.get("title")),
            ("description", properties.get("descr")),
        )
        explicit = [(source, value) for source, value in values if value]
        native = _native_locator(slide_number, shape_path)
        if not explicit:
            warnings.append(_warning("PPTX_IMAGE_WITHOUT_ALT", native))
            return
        if not preserve_image_alt:
            return
        for alt_source, text in explicit:
            builder.add(
                "image_alt",
                text,
                _native_locator(
                    slide_number,
                    shape_path,
                    alt_source=alt_source,
                ),
                {"alt_source": alt_source},
                paragraph_like=False,
            )

    @staticmethod
    def _visit_notes(
        slide,
        slide_number: int,
        builder: _DocumentBuilder,
        warnings: list[dict[str, object]],
        preserve_notes: bool,
    ) -> None:
        if not slide.has_notes_slide:
            return
        text_frame = slide.notes_slide.notes_text_frame
        if not any(normalize_quote(paragraph.text) for paragraph in text_frame.paragraphs):
            return
        note_shape = next(
            shape
            for shape in slide.notes_slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text_frame._txBody is text_frame._txBody
        )
        tree_index = next(
            index
            for index, shape in enumerate(slide.notes_slide.shapes)
            if shape._element is note_shape._element
        )
        shape_path = ({"shape_id": note_shape.shape_id, "tree_index": tree_index},)
        if not preserve_notes:
            warnings.append(
                _warning(
                    "PPTX_NOTES_SKIPPED",
                    _native_locator(slide_number, shape_path, notes=True),
                )
            )
            return
        for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
            builder.add(
                "footnote",
                paragraph.text,
                _native_locator(
                    slide_number,
                    shape_path,
                    text_frame_paragraph_index=paragraph_index,
                    notes=True,
                ),
                {},
                paragraph_like=True,
            )
